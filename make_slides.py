# -*- coding: utf-8 -*-
"""
Build slides.pptx from IVE Identity-Aware Segmentation project.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BASE = r"E:\Master_Degree\DL-FOR-COMPUTER-VISION\Segmentation_Ive"
OUT  = os.path.join(BASE, "slides.pptx")
RES  = os.path.join(BASE, "Result")

# --- Colors ---
C_BG     = RGBColor(0x0F, 0x0F, 0x1A)
C_CARD   = RGBColor(0x1A, 0x1A, 0x30)
C_PINK   = RGBColor(0xFF, 0x6B, 0x9D)
C_PURPLE = RGBColor(0x7B, 0x61, 0xFF)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xCC, 0xCC, 0xDD)
C_YELLOW = RGBColor(0xFF, 0xE0, 0x66)
C_GREEN  = RGBColor(0x66, 0xFF, 0xB2)
C_RED    = RGBColor(0xFF, 0x66, 0x66)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# --- Helpers ---
def add_bg(slide, color=C_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def txb(slide, text, left, top, width, height,
        size=20, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def mtxb(slide, lines, left, top, width, height, size=18, color=C_LIGHT, gap=3):
    """multi-line textbox - each item in list is a paragraph"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(gap)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb

def box(slide, left, top, width, height, fill=C_CARD):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def bar(slide, left, top, width, height, color):
    b = box(slide, left, top, width, height, fill=color)
    return b

def hbar(slide):
    bar(slide, 0, 0, W, Inches(0.07), C_PINK)

def tag(slide, label):
    box(slide, Inches(0.4), Inches(0.28), Inches(2.5), Inches(0.38), fill=C_PURPLE)
    txb(slide, label, Inches(0.4), Inches(0.28), Inches(2.5), Inches(0.38),
        size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

def h1(slide, text, top=Inches(0.82)):
    txb(slide, text, Inches(0.4), top, Inches(12.5), Inches(0.7),
        size=32, bold=True, color=C_WHITE)

def add_img(slide, path, left, top, width=None, height=None):
    if width and height:
        slide.shapes.add_picture(path, left, top, width, height)
    elif width:
        slide.shapes.add_picture(path, left, top, width=width)
    elif height:
        slide.shapes.add_picture(path, left, top, height=height)
    else:
        slide.shapes.add_picture(path, left, top)

# ================================================================
# S1 - TITLE
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
bar(s, 0, 0, W, Inches(0.07), C_PINK)
bar(s, 0, H - Inches(0.07), W, Inches(0.07), C_PURPLE)

# deco circles
for cx, cy, r, col in [
    (Inches(9.5), Inches(-0.8), Inches(5.0), RGBColor(0x1E, 0x0A, 0x3C)),
    (Inches(10.8), Inches(3.8), Inches(3.5), RGBColor(0x28, 0x08, 0x40)),
]:
    sh = s.shapes.add_shape(9, cx, cy, r, r)
    sh.fill.solid(); sh.fill.fore_color.rgb = col
    sh.line.fill.background()

txb(s, "IVE", Inches(0.7), Inches(0.9), Inches(10), Inches(1.6),
    size=80, bold=True, color=C_PINK)
txb(s, "Identity-Aware Text-Prompt Segmentation",
    Inches(0.7), Inches(2.4), Inches(11), Inches(1.0),
    size=32, bold=True, color=C_WHITE)
txb(s, "InsightFace  x  SAM3  Pipeline",
    Inches(0.7), Inches(3.35), Inches(10), Inches(0.6),
    size=24, color=C_PURPLE)
bar(s, Inches(0.7), Inches(4.05), Inches(5.5), Inches(0.04), C_PINK)
txb(s, "Deep Learning for Computer Vision  |  2026",
    Inches(0.7), Inches(4.2), Inches(10), Inches(0.5),
    size=17, color=C_LIGHT)

# ================================================================
# S2 - OVERVIEW
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s); hbar(s); tag(s, "OVERVIEW")
h1(s, "ระบบทำอะไร?")

txb(s,
    "ระบบ segmentation ที่รู้จักตัวตน (Identity-Aware)\n"
    "ไม่ใช่แค่หา 'person' แต่รู้ว่าคนนั้นคือใคร และ segment เฉพาะส่วนที่เป็นของเขา",
    Inches(0.5), Inches(1.75), Inches(12.3), Inches(1.1),
    size=22, color=C_LIGHT)

cards = [
    (C_PINK,   "Identity\nRecognition",  "แยกสมาชิก IVE 6 คน\nด้วยใบหน้า ArcFace"),
    (C_PURPLE, "Possession\nAware",      '"yujin shirt" = เสื้อ\nของยูจินเท่านั้น'),
    (C_GREEN,  "Text Prompt\nDriven",    "พิมพ์ชื่อ + วัตถุ\nระบบทำให้เอง"),
]
for i, (col, title, desc) in enumerate(cards):
    cx = Inches(0.5 + i * 4.25)
    cy = Inches(3.0)
    box(s, cx, cy, Inches(3.9), Inches(3.3))
    bar(s, cx, cy, Inches(3.9), Inches(0.07), col)
    txb(s, title, cx + Inches(0.2), cy + Inches(0.25),
        Inches(3.6), Inches(1.0), size=22, bold=True, color=col)
    txb(s, desc, cx + Inches(0.2), cy + Inches(1.25),
        Inches(3.6), Inches(1.8), size=18, color=C_LIGHT)

# ================================================================
# S3 - HIDDEN AGENDA
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s); hbar(s); tag(s, "OBJECTIVES")
h1(s, "3 Hidden Agenda ของโจทย์")

agendas = [
    (C_PINK,   "01", "Object -> Identity",
     "SAM3 มาตรฐานมองทุกคนเป็น 'person' เหมือนกันหมด\nโจทย์ต้องการแยก wonyoung / yujin / gaeul ในภาพหมู่"),
    (C_PURPLE, "02", "Possession Aware",
     '"yujin shirt" ต้องหมายถึงเสื้อของยูจินเท่านั้น\nไม่ใช่เสื้อของคนที่ยืนข้างๆ'),
    (C_YELLOW, "03", "Cross-lingual Prompt",
     'ระบบควรรองรับ prompt ภาษาอื่น เช่น ภาษาไทย\n"ขอเสื้อของน้องวอนยอง"'),
]
for i, (col, num, title, desc) in enumerate(agendas):
    cy = Inches(1.75 + i * 1.75)
    box(s, Inches(0.5), cy, Inches(0.95), Inches(1.35), fill=C_CARD)
    txb(s, num, Inches(0.5), cy + Inches(0.18), Inches(0.95), Inches(0.9),
        size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(s, title, Inches(1.65), cy + Inches(0.08),
        Inches(11), Inches(0.5), size=23, bold=True, color=col)
    txb(s, desc, Inches(1.65), cy + Inches(0.58),
        Inches(11.2), Inches(0.8), size=18, color=C_LIGHT)

# ================================================================
# S4 - FIRST APPROACH (FINE-TUNING)
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s); hbar(s); tag(s, "APPROACH #1")
h1(s, "แนวทางแรก: Fine-tuning SAM3 โดยตรง")

# left panel
box(s, Inches(0.4), Inches(1.6), Inches(5.9), Inches(5.2))
txb(s, "สิ่งที่ทำ", Inches(0.6), Inches(1.75), Inches(5.6), Inches(0.5),
    size=20, bold=True, color=C_PURPLE)
mtxb(s, [
    "->  เทรน SAM3 ด้วย dataset รูปเดี่ยวสมาชิก IVE",
    "->  ใส่ text label ชื่อบุคคลแต่ละคน",
    "->  หวังให้โมเดลเรียนรู้ความแตกต่างของคน",
], Inches(0.6), Inches(2.3), Inches(5.6), Inches(2.5), size=18, gap=5)

# right panel
box(s, Inches(6.9), Inches(1.6), Inches(6.0), Inches(5.2),
    fill=RGBColor(0x25, 0x10, 0x18))
txb(s, "ปัญหาที่พบ  X", Inches(7.1), Inches(1.75), Inches(5.7), Inches(0.5),
    size=20, bold=True, color=C_RED)
mtxb(s, [
    "X  Dataset ไม่เพียงพอ — ต้องการข้อมูลมหาศาล",
    "X  Model bias — SAM3 ยังมองทุกคนเป็น 'person'",
    "X  ใช้กับภาพหมู่ไม่ได้ — แยกแยะตัวตนไม่ออก",
], Inches(7.1), Inches(2.3), Inches(5.7), Inches(2.5),
    size=18, color=RGBColor(0xFF, 0xAA, 0xAA), gap=5)

txb(s, "->  Fine-tuning = Dead End  ->  เปลี่ยนแนวทาง",
    Inches(1.5), Inches(6.2), Inches(10), Inches(0.55),
    size=22, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)

# ================================================================
# S5 - PIPELINE ARCHITECTURE
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s); hbar(s); tag(s, "SOLUTION")
h1(s, "InsightFace + SAM3 Pipeline")

steps = [
    (C_PINK,   'Text Prompt\n"wonyoung shirt"'),
    (C_PURPLE, 'Prompt Parser\nmember + SAM3 text'),
    (C_GREEN,  'InsightFace (ArcFace)\nFace detect + embed'),
    (C_YELLOW, 'Cosine Similarity\nMatch DB 6 members'),
    (C_PINK,   'BBox Expansion\nFace -> Body region'),
    (C_PURPLE, 'SAM3 Segmentation\ntext-prompted on ROI'),
    (C_GREEN,  'Mask Scoring\nconfidence + proximity'),
    (C_YELLOW, 'Final Output\nMask on full image'),
]

bw = Inches(1.38)
bh = Inches(1.12)
gx = Inches(0.2)
sx = Inches(0.22)
sy = Inches(1.72)
row_gap = Inches(1.18)

for i, (col, label) in enumerate(steps):
    r = i // 4
    c = i % 4
    bx = sx + c * (bw + gx)
    by = sy + r * (bh + row_gap)
    box(s, bx, by, bw, bh)
    bar(s, bx, by, bw, Inches(0.06), col)
    txb(s, label, bx + Inches(0.08), by + Inches(0.12),
        bw - Inches(0.1), bh - Inches(0.15), size=14, color=C_WHITE)
    if c < 3:
        txb(s, "->", bx + bw + Inches(0.02), by + Inches(0.35),
            gx, Inches(0.4), size=18, bold=True, color=col, align=PP_ALIGN.CENTER)

# down arrow
txb(s, "v  (ต่อแถวล่าง)", Inches(6.0), sy + bh + Inches(0.08),
    Inches(2.0), Inches(0.5), size=14, color=C_LIGHT)

# ================================================================
# S6 - PIPELINE DETAILS
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s); hbar(s); tag(s, "DETAILS")
h1(s, "รายละเอียดแต่ละขั้นตอน")

details = [
    (C_PINK,   "Prompt Parser",
     '"wonyoung shirt"  ->  member="Wonyoung"  +  SAM3_text="shirt"'),
    (C_PURPLE, "InsightFace (ArcFace)",
     "RetinaFace detect -> ArcFace embed -> Cosine similarity vs DB\nThreshold: 0.42"),
    (C_GREEN,  "BBox Expansion",
     "กว้าง x3.0  |  บน x1.2 (รวมผม)  |  ล่าง x5.0 (ครอบร่างกาย)"),
    (C_YELLOW, "Mask Quality Scoring",
     "composite = confidence + proximity + overlap\nproximity = centroid ของ mask ใกล้ center ROI เท่าไหร่"),
]
for i, (col, title, desc) in enumerate(details):
    r = i // 2
    c = i % 2
    bx = Inches(0.4 + c * 6.5)
    by = Inches(1.7 + r * 2.55)
    bw = Inches(6.15)
    bh = Inches(2.3)
    box(s, bx, by, bw, bh)
    bar(s, bx, by, Inches(0.07), bh, col)
    txb(s, title, bx + Inches(0.2), by + Inches(0.15),
        bw - Inches(0.3), Inches(0.5), size=20, bold=True, color=col)
    txb(s, desc, bx + Inches(0.2), by + Inches(0.68),
        bw - Inches(0.3), Inches(1.45), size=17, color=C_LIGHT)

# ================================================================
# S7 - PROBLEM & SOLUTION
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s); hbar(s); tag(s, "CHALLENGE")
h1(s, "ปัญหา BBox ทับซ้อน & วิธีแก้")

box(s, Inches(0.4), Inches(1.65), Inches(12.5), Inches(2.35),
    fill=RGBColor(0x25, 0x10, 0x18))
bar(s, Inches(0.4), Inches(1.65), Inches(0.07), Inches(2.35), C_RED)
txb(s, "ปัญหา  X", Inches(0.65), Inches(1.75), Inches(12), Inches(0.5),
    size=21, bold=True, color=C_RED)
txb(s,
    'BBox ที่ขยายแล้วอาจครอบคนข้างๆด้วย\n'
    '->  prompt "wonyoung hair"  แต่ SAM3 segment ผมของคนข้างๆแทน',
    Inches(0.65), Inches(2.32), Inches(12), Inches(1.1),
    size=19, color=C_LIGHT)

box(s, Inches(0.4), Inches(4.15), Inches(12.5), Inches(2.9),
    fill=RGBColor(0x10, 0x25, 0x18))
bar(s, Inches(0.4), Inches(4.15), Inches(0.07), Inches(2.9), C_GREEN)
txb(s, "วิธีแก้: Proximity Score  OK", Inches(0.65), Inches(4.25),
    Inches(12), Inches(0.5), size=21, bold=True, color=C_GREEN)
txb(s,
    "คำนวณ centroid ของแต่ละ mask candidate\n"
    "->  mask ที่อยู่ใกล้ center ของ ROI จะได้ composite score สูงกว่า\n"
    "->  เลือก mask ที่ถูกต้องได้แม้ BBox ทับซ้อนกัน",
    Inches(0.65), Inches(4.8), Inches(12), Inches(1.8),
    size=19, color=C_LIGHT)

box(s, Inches(3.0), Inches(3.55), Inches(7.3), Inches(0.52),
    fill=RGBColor(0x20, 0x20, 0x40))
txb(s, "composite = confidence + proximity + overlap",
    Inches(3.15), Inches(3.6), Inches(7.0), Inches(0.45),
    size=19, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)

# ================================================================
# S8 - RESULT: WONYOUNG
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s); hbar(s); tag(s, "RESULTS")
h1(s, 'Prompt: "wonyoung"')

img_path = os.path.join(RES, "Wonyoung.png")
img_h = Inches(5.15)
img_w = Inches(5.15 * 1956 / 895)
img_x = (W - img_w) / 2
add_img(s, img_path, img_x, Inches(1.7), width=img_w, height=img_h)

stats = [("SAM3 text", '"person"'), ("Similarity", "0.664"), ("Pixels", "316,909")]
for i, (k, v) in enumerate(stats):
    bx = Inches(0.4 + i * 4.15)
    box(s, bx, Inches(6.88), Inches(3.9), Inches(0.48))
    txb(s, k + ": " + v, bx + Inches(0.15), Inches(6.93),
        Inches(3.6), Inches(0.38), size=18, bold=True, color=C_PINK)

# ================================================================
# S9 - RESULT: YUJIN
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s); hbar(s); tag(s, "RESULTS")
h1(s, '"yujin face"  &  "yujin shirt"')

face_path  = os.path.join(RES, "Yujin_face.png")
shirt_path = os.path.join(RES, "Yujin_shirt.png")

face_h  = Inches(2.85)
face_w  = Inches(2.85 * 2010 / 905)
shirt_h = Inches(2.85)
shirt_w = Inches(2.85 * 1938 / 905)

add_img(s, face_path,  Inches(0.1),  Inches(1.8), width=face_w,  height=face_h)
add_img(s, shirt_path, Inches(6.8),  Inches(1.8), width=shirt_w, height=shirt_h)

txb(s, 'Prompt: "yujin face"', Inches(0.1), Inches(4.85),
    Inches(6.0), Inches(0.5), size=19, bold=True, color=C_PINK)
mtxb(s, ['SAM3 text: "face"', 'Similarity: 0.663', 'Pixels: 31,049'],
     Inches(0.1), Inches(5.3), Inches(6.0), Inches(1.5), size=17, color=C_LIGHT)

txb(s, 'Prompt: "yujin shirt"', Inches(6.8), Inches(4.85),
    Inches(6.0), Inches(0.5), size=19, bold=True, color=C_PURPLE)
mtxb(s, ['SAM3 text: "shirt"', 'Similarity: 0.663', 'Pixels: 244,312'],
     Inches(6.8), Inches(5.3), Inches(6.0), Inches(1.5), size=17, color=C_LIGHT)

box(s, Inches(3.6), Inches(3.0), Inches(6.1), Inches(0.6),
    fill=RGBColor(0x10, 0x25, 0x18))
txb(s, "OK  ได้เสื้อของยูจิน ไม่ใช่เสื้อของคนข้างๆ  (Possession-aware)",
    Inches(3.7), Inches(3.06), Inches(5.9), Inches(0.48),
    size=17, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

# ================================================================
# S10 - VIDEO RESULTS
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s); hbar(s); tag(s, "RESULTS")
h1(s, "Video Results")

videos = [
    ("all_members_video.mp4",
     "All 6 Members — Segment + BBox ทุกสมาชิก IVE"),
    ("all_members_video_remove_bbox.mp4",
     "All 6 Members — เวอร์ชันไม่มี BBox"),
    ("test_wonyoung.mp4",
     'Prompt "wonyoung" — 909 frames, 78.9% detection rate'),
    ("test_wonyoung_hair.mp4",
     'Prompt "wonyoung hair" — segment ผมวอนยอง'),
    ("yujin_hair.mp4",
     'Prompt "yujin hair" — segment ผมยูจิน'),
]
for i, (fname, desc) in enumerate(videos):
    cy = Inches(1.7 + i * 1.02)
    box(s, Inches(0.4), cy, Inches(12.5), Inches(0.88))
    txb(s, "[vid]  " + fname, Inches(0.6), cy + Inches(0.07),
        Inches(5.5), Inches(0.4), size=17, bold=True, color=C_PINK)
    txb(s, desc, Inches(6.3), cy + Inches(0.07),
        Inches(6.5), Inches(0.4), size=17, color=C_LIGHT)

# ================================================================
# S11 - VIDEO STATS
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s); hbar(s); tag(s, "STATS")
h1(s, "สถิติการประมวลผลวิดีโอ")

rows = [
    ("Input Video",        "IVE-30s.mp4"),
    ("ความละเอียด",        "1920 x 1080"),
    ("Frame rate",         "30 fps"),
    ("Frames ทั้งหมด",     "909 frames"),
    ("Frame sampling",     "ทุก 3 frames"),
    ("Frames ที่ประมวลผล", "303 frames"),
    ("Detection success",  "239 / 303  (78.9%)"),
    ("เวลาประมวลผล",       "~82.8 วินาที  (~11 fps effective)"),
]
alt = [RGBColor(0x1E, 0x1E, 0x38), RGBColor(0x18, 0x18, 0x2E)]
bar(s, Inches(0.4), Inches(1.58), Inches(5.5), Inches(0.1), C_PINK)
bar(s, Inches(5.95), Inches(1.58), Inches(7.0), Inches(0.1), C_PURPLE)
for i, (k, v) in enumerate(rows):
    cy = Inches(1.68 + i * 0.67)
    bg = alt[i % 2]
    box(s, Inches(0.4), cy, Inches(5.5), Inches(0.62), fill=bg)
    box(s, Inches(5.95), cy, Inches(7.0), Inches(0.62), fill=bg)
    txb(s, k, Inches(0.6), cy + Inches(0.1), Inches(5.2), Inches(0.45), size=19, color=C_LIGHT)
    hi = C_GREEN if i == 6 else C_WHITE
    txb(s, v, Inches(6.1), cy + Inches(0.1), Inches(6.7), Inches(0.45),
        size=19, bold=(i == 6), color=hi)

# ================================================================
# S12 - Q1: THAI PROMPT
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s); hbar(s); tag(s, "Q&A")
h1(s, "Q1: รองรับ Thai Prompt ได้ไหม?")

txb(s, "ปัจจุบัน: ไม่รองรับ", Inches(0.5), Inches(1.68),
    Inches(12), Inches(0.5), size=21, bold=True, color=C_RED)
txb(s, "SAM3 text encoder เทรนด้วย dataset ภาษาอังกฤษ -> ไม่เข้าใจ semantic ภาษาไทย",
    Inches(0.5), Inches(2.22), Inches(12), Inches(0.5), size=18, color=C_LIGHT)

# panel left - long term
box(s, Inches(0.4), Inches(2.92), Inches(5.9), Inches(3.9))
bar(s, Inches(0.4), Inches(2.92), Inches(5.9), Inches(0.07), C_PURPLE)
txb(s, "ระยะยาว", Inches(0.6), Inches(3.02), Inches(5.6), Inches(0.5),
    size=20, bold=True, color=C_PURPLE)
txb(s,
    "Fine-tune SAM3 text encoder\nด้วย multilingual dataset\nที่มี caption + text prompt ภาษาไทย",
    Inches(0.6), Inches(3.6), Inches(5.6), Inches(2.0), size=18, color=C_LIGHT)

# panel right - short term
box(s, Inches(7.0), Inches(2.92), Inches(5.9), Inches(3.9),
    fill=RGBColor(0x10, 0x25, 0x18))
bar(s, Inches(7.0), Inches(2.92), Inches(5.9), Inches(0.07), C_GREEN)
txb(s, "ระยะสั้น (Practical)  OK", Inches(7.2), Inches(3.02),
    Inches(5.6), Inches(0.5), size=20, bold=True, color=C_GREEN)
box(s, Inches(7.2), Inches(3.6), Inches(5.5), Inches(1.85),
    fill=RGBColor(0x1E, 0x1E, 0x3A))
flow = (
    '"ขอเสื้อของน้องวอนยอง"\n'
    '          v  LLM\n'
    '     "wonyoung shirt"\n'
    '          v  Pipeline'
)
txb(s, flow, Inches(7.35), Inches(3.68), Inches(5.2), Inches(1.7),
    size=16, color=RGBColor(0xA0, 0xE0, 0xFF))
txb(s, "ไม่ต้อง fine-tune ใหม่ | LLM เข้าใจ context ได้ดี",
    Inches(7.2), Inches(5.55), Inches(5.6), Inches(0.9),
    size=16, color=C_LIGHT)

# ================================================================
# S13 - Q2: POSSESSION
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s); hbar(s); tag(s, "Q&A")
h1(s, 'Q2: "yujin shirt" = เสื้อของยูจินเท่านั้น')

steps4 = [
    (C_PINK,   "1", "InsightFace\nหายูจิน\n-> face bbox"),
    (C_PURPLE, "2", "ขยาย BBox\n-> ROI ครอบ\nร่างกายยูจิน"),
    (C_GREEN,  "3", 'SAM3 หา\n"shirt" ภายใน\nROI เท่านั้น'),
    (C_YELLOW, "4", "ผลลัพธ์ =\nเสื้อของยูจิน\nตามคำนิยาม"),
]
bw4 = Inches(2.8)
bh4 = Inches(3.6)
gx4 = Inches(0.55)
sx4 = Inches(0.7)

for i, (col, num, desc) in enumerate(steps4):
    bx = sx4 + i * (bw4 + gx4)
    by = Inches(2.0)
    box(s, bx, by, bw4, bh4)
    bar(s, bx, by, bw4, Inches(0.07), col)
    txb(s, num, bx, by + Inches(0.15), bw4, Inches(0.9),
        size=48, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(s, desc, bx + Inches(0.15), by + Inches(1.1),
        bw4 - Inches(0.2), Inches(2.2), size=17, color=C_WHITE,
        align=PP_ALIGN.CENTER)
    if i < 3:
        txb(s, "->", bx + bw4 + Inches(0.1), by + Inches(1.55),
            Inches(0.4), Inches(0.7), size=22, bold=True, color=col)

box(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.62),
    fill=RGBColor(0x10, 0x25, 0x18))
txb(s, "Key Insight: จำกัด search space ให้เหลือแค่ ROI ของบุคคลที่ระบุ",
    Inches(0.5), Inches(5.92), Inches(12.3), Inches(0.5),
    size=20, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

# ================================================================
# S14 - Q3: CROSS-LINGUAL
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s); hbar(s); tag(s, "Q&A")
h1(s, "Q3: Cross-lingual — สรุปแนวทาง")

# table headers
for lx, lw, label, col in [
    (0.4, 4.3, "แนวทาง",  C_PINK),
    (4.75, 3.8, "ข้อดี",   C_GREEN),
    (8.6, 4.3, "ข้อเสีย", C_RED),
]:
    bar(s, Inches(lx), Inches(1.68), Inches(lw), Inches(0.55), col)
    txb(s, label, Inches(lx + 0.15), Inches(1.73), Inches(lw), Inches(0.45),
        size=18, bold=True, color=C_WHITE)

table_data = [
    ("Fine-tune SAM3\nด้วย Thai data",
     "รองรับภาษาไทย\nโดยตรง",
     "Dataset ขนาดใหญ่\nใช้เวลานาน"),
    ("LLM translation layer\n(แนะนำ)",
     "ทำได้ทันที\nไม่ต้อง fine-tune",
     "ขึ้น external API\nlatency เพิ่มขึ้น"),
]
alt2 = [RGBColor(0x1A, 0x1A, 0x30), RGBColor(0x14, 0x22, 0x14)]
for ri, (a, b, c) in enumerate(table_data):
    cy = Inches(2.3 + ri * 1.6)
    bg = alt2[ri]
    for lx, lw, text, tc in [
        (0.4,  4.3, a, C_LIGHT),
        (4.75, 3.8, b, C_GREEN if ri == 1 else C_LIGHT),
        (8.6,  4.3, c, C_LIGHT),
    ]:
        box(s, Inches(lx), cy, Inches(lw), Inches(1.5), fill=bg)
        txb(s, text, Inches(lx + 0.15), cy + Inches(0.2),
            Inches(lw - 0.2), Inches(1.15), size=17, color=tc)

box(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.75),
    fill=RGBColor(0x20, 0x20, 0x10))
txb(s,
    "สรุป: LLM translation layer practical กว่า\n"
    "LLM สมัยใหม่รองรับหลายภาษาได้ดีอยู่แล้ว",
    Inches(0.6), Inches(5.65), Inches(12.1), Inches(0.65),
    size=18, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)

# ================================================================
# S15 - SUMMARY
# ================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
bar(s, 0, 0, W, Inches(0.07), C_PINK)
bar(s, 0, H - Inches(0.07), W, Inches(0.07), C_PURPLE)

txb(s, "สรุป", Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.85),
    size=44, bold=True, color=C_PINK, align=PP_ALIGN.CENTER)

summary = [
    (C_PINK,   '"person"  ->  "wonyoung"',
     "Identity-aware: แยกตัวตนด้วย ArcFace face embedding"),
    (C_PURPLE, '"shirt in image"  ->  "shirt ของยูจิน"',
     "Possession-aware: Segment ใน ROI เฉพาะบุคคล"),
    (C_GREEN,  "Pipeline Modular",
     "เพิ่มสมาชิกใหม่ได้ทันที -> เพิ่ม face embedding ใน DB"),
    (C_YELLOW, "Cross-lingual Support",
     "Short-term: LLM translation layer แปลง Thai -> EN prompt"),
]
for i, (col, title, desc) in enumerate(summary):
    cy = Inches(1.42 + i * 1.42)
    box(s, Inches(0.4), cy, Inches(12.5), Inches(1.28))
    bar(s, Inches(0.4), cy, Inches(0.08), Inches(1.28), col)
    txb(s, title, Inches(0.65), cy + Inches(0.1),
        Inches(5.5), Inches(0.5), size=20, bold=True, color=col)
    txb(s, desc, Inches(0.65), cy + Inches(0.62),
        Inches(11.8), Inches(0.55), size=18, color=C_LIGHT)

txb(s, "InsightFace  x  SAM3  =  Identity-Aware Segmentation",
    Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.38),
    size=17, bold=True, color=C_PURPLE, align=PP_ALIGN.CENTER)

# ================================================================
# SAVE
# ================================================================
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
